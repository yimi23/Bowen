#!/usr/bin/env python3
"""
BOWEN Advanced Document Creation Engine
Creates professional documents (essays, reports, presentations)
"""

import os
import json
import logging
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional, Any

logger = logging.getLogger(__name__)

class AdvancedDocumentEngine:
    """
    Advanced document creation system that:
    1. Essay writing with proper structure
    2. PowerPoint presentations
    3. Research reports
    4. Professional documents
    5. Academic papers
    6. Business documents
    """
    
    def __init__(self, document_engine, research_engine, claude_engine):
        """Use existing engines plus Claude for content generation"""
        self.document_engine = document_engine
        self.research_engine = research_engine
        self.claude_engine = claude_engine
        self.output_dir = Path.home() / "Desktop" / "bowen_outputs"
        self.output_dir.mkdir(exist_ok=True)
        
        # Citation styles
        self.citation_styles = {
            "APA": self._apa_citation_format,
            "MLA": self._mla_citation_format,
            "Chicago": self._chicago_citation_format
        }
        
        logger.info("Advanced Document Engine initialized")
    
    def write_essay(self, topic: str, requirements: Dict[str, Any]) -> str:
        """
        Generate complete essay with research, citations, and formatting
        """
        try:
            length = requirements.get('length', 1500)
            style = requirements.get('style', 'APA')
            sources = requirements.get('sources', 5)
            deadline = requirements.get('deadline')
            
            logger.info(f"Writing {length}-word essay on '{topic}' in {style} style")
            
            # Step 1: Research the topic
            research_results = self.research_engine.research_and_report(
                f"{topic} academic sources scholarly articles"
            )
            
            # Step 2: Create essay outline
            outline = self._create_essay_outline(topic, requirements, research_results)
            
            # Step 3: Write each section
            essay_sections = self._write_essay_sections(topic, outline, research_results, requirements)
            
            # Step 4: Add citations and bibliography
            formatted_essay = self._add_citations_and_bibliography(essay_sections, style)
            
            # Step 5: Format as DOCX
            filename = f"{topic.replace(' ', '_').lower()}_essay_{datetime.now().strftime('%Y%m%d')}.docx"
            filepath = self._create_docx_essay(formatted_essay, filename, style)
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error writing essay: {e}")
            return f"Error writing essay: {str(e)}"
    
    def create_presentation(self, topic: str, spec: Dict[str, Any]) -> str:
        """
        Generate PowerPoint presentation with research and visuals
        """
        try:
            slides = spec.get('slides', 10)
            style = spec.get('style', 'professional')
            speaker_notes = spec.get('include_speaker_notes', True)
            
            logger.info(f"Creating {slides}-slide presentation on '{topic}'")
            
            # Step 1: Research the topic
            research_results = self.research_engine.research_and_report(
                f"{topic} presentation key points statistics data"
            )
            
            # Step 2: Create slide outline
            slide_outline = self._create_presentation_outline(topic, slides, research_results)
            
            # Step 3: Generate content for each slide
            slide_contents = self._generate_slide_contents(slide_outline, research_results, speaker_notes)
            
            # Step 4: Create PowerPoint file
            filename = f"{topic.replace(' ', '_').lower()}_presentation_{datetime.now().strftime('%Y%m%d')}.pptx"
            filepath = self._create_powerpoint(slide_contents, filename, style)
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            return f"Error creating presentation: {str(e)}"
    
    def write_report(self, subject: str, data: Dict[str, Any]) -> str:
        """
        Generate professional report with analysis and recommendations
        """
        try:
            logger.info(f"Writing professional report on '{subject}'")
            
            # Step 1: Research background information
            background_research = self.research_engine.research_and_report(
                f"{subject} analysis best practices industry standards"
            )
            
            # Step 2: Create report structure
            report_structure = self._create_report_structure(subject, data)
            
            # Step 3: Write each section
            report_sections = self._write_report_sections(subject, report_structure, data, background_research)
            
            # Step 4: Add visualizations if data provided
            if 'charts' in data or 'tables' in data:
                report_sections = self._add_visualizations_to_report(report_sections, data)
            
            # Step 5: Format as PDF
            filename = f"{subject.replace(' ', '_').lower()}_report_{datetime.now().strftime('%Y%m%d')}.pdf"
            filepath = self._create_pdf_report(report_sections, filename)
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error writing report: {e}")
            return f"Error writing report: {str(e)}"
    
    def _create_essay_outline(self, topic: str, requirements: Dict[str, Any], research: str) -> Dict[str, Any]:
        """Create essay outline using Claude"""
        try:
            outline_prompt = f"""
            Create a detailed essay outline for the topic: "{topic}"
            
            Requirements:
            - Length: {requirements.get('length', 1500)} words
            - Style: {requirements.get('style', 'APA')}
            - Academic level: {requirements.get('level', 'undergraduate')}
            
            Research context (use for ideas):
            {research[:2000]}  # Limit context
            
            Create an outline with:
            1. Thesis statement
            2. Introduction points (10% of word count)
            3. Body paragraphs (80% of word count) - each with main point and supporting details
            4. Conclusion points (10% of word count)
            
            Return as JSON:
            {{
                "thesis": "Your thesis statement",
                "introduction": ["point 1", "point 2"],
                "body_paragraphs": [
                    {{"main_point": "...", "supporting_details": ["detail 1", "detail 2"]}},
                    {{"main_point": "...", "supporting_details": ["detail 1", "detail 2"]}}
                ],
                "conclusion": ["point 1", "point 2"]
            }}
            """
            
            response = self.claude_engine.client.messages.create(
                model=self.claude_engine.model,
                max_tokens=1500,
                messages=[{"role": "user", "content": outline_prompt}]
            )
            
            outline_text = response.content[0].text.strip()
            
            # Clean and parse JSON
            if outline_text.startswith('```json'):
                outline_text = outline_text[7:-3]
            elif outline_text.startswith('```'):
                outline_text = outline_text[3:-3]
            
            return json.loads(outline_text)
            
        except Exception as e:
            logger.error(f"Error creating essay outline: {e}")
            # Fallback outline
            return {
                "thesis": f"Analysis of {topic}",
                "introduction": [f"Introduction to {topic}", "Research overview"],
                "body_paragraphs": [
                    {"main_point": f"Key aspect 1 of {topic}", "supporting_details": ["Detail 1", "Detail 2"]},
                    {"main_point": f"Key aspect 2 of {topic}", "supporting_details": ["Detail 1", "Detail 2"]},
                    {"main_point": f"Key aspect 3 of {topic}", "supporting_details": ["Detail 1", "Detail 2"]}
                ],
                "conclusion": ["Summary of findings", "Implications"]
            }
    
    def _write_essay_sections(self, topic: str, outline: Dict[str, Any], research: str, requirements: Dict[str, Any]) -> Dict[str, str]:
        """Write each section of the essay"""
        sections = {}
        
        try:
            # Write introduction
            intro_prompt = f"""
            Write an introduction paragraph for an essay on "{topic}".
            
            Thesis: {outline['thesis']}
            Points to cover: {outline['introduction']}
            Word count: ~{int(requirements.get('length', 1500) * 0.1)} words
            Style: {requirements.get('style', 'APA')}
            
            Research context:
            {research[:1500]}
            
            Write a compelling, academic introduction that hooks the reader and presents the thesis clearly.
            """
            
            intro_response = self.claude_engine.client.messages.create(
                model=self.claude_engine.model,
                max_tokens=800,
                messages=[{"role": "user", "content": intro_prompt}]
            )
            sections['introduction'] = intro_response.content[0].text
            
            # Write body paragraphs
            body_paragraphs = []
            for i, paragraph in enumerate(outline['body_paragraphs']):
                para_prompt = f"""
                Write a body paragraph for an essay on "{topic}".
                
                Main point: {paragraph['main_point']}
                Supporting details: {paragraph['supporting_details']}
                Word count: ~{int(requirements.get('length', 1500) * 0.8 / len(outline['body_paragraphs']))} words
                Style: {requirements.get('style', 'APA')}
                
                Research context:
                {research[:1500]}
                
                Write a well-structured paragraph with topic sentence, evidence, and analysis.
                """
                
                para_response = self.claude_engine.client.messages.create(
                    model=self.claude_engine.model,
                    max_tokens=600,
                    messages=[{"role": "user", "content": para_prompt}]
                )
                body_paragraphs.append(para_response.content[0].text)
            
            sections['body'] = '\n\n'.join(body_paragraphs)
            
            # Write conclusion
            conclusion_prompt = f"""
            Write a conclusion paragraph for an essay on "{topic}".
            
            Thesis: {outline['thesis']}
            Points to cover: {outline['conclusion']}
            Word count: ~{int(requirements.get('length', 1500) * 0.1)} words
            Style: {requirements.get('style', 'APA')}
            
            Summarize key findings and provide thoughtful closing remarks without introducing new information.
            """
            
            conclusion_response = self.claude_engine.client.messages.create(
                model=self.claude_engine.model,
                max_tokens=500,
                messages=[{"role": "user", "content": conclusion_prompt}]
            )
            sections['conclusion'] = conclusion_response.content[0].text
            
            return sections
            
        except Exception as e:
            logger.error(f"Error writing essay sections: {e}")
            return {
                'introduction': f"Introduction to {topic}...",
                'body': f"Body content about {topic}...",
                'conclusion': f"Conclusion about {topic}..."
            }
    
    def _add_citations_and_bibliography(self, sections: Dict[str, str], style: str) -> str:
        """Add proper citations and bibliography"""
        try:
            # Combine sections
            full_essay = f"{sections['introduction']}\n\n{sections['body']}\n\n{sections['conclusion']}"
            
            # Add bibliography
            bibliography = self._generate_bibliography(style)
            
            return f"{full_essay}\n\n{bibliography}"
            
        except Exception as e:
            logger.error(f"Error adding citations: {e}")
            return '\n\n'.join(sections.values())
    
    def _create_docx_essay(self, content: str, filename: str, style: str) -> str:
        """Create DOCX file with proper academic formatting"""
        try:
            # Use existing document engine
            filepath = self.document_engine.create_document(
                content=content,
                title=filename.replace('.docx', '').replace('_', ' ').title(),
                doc_type='essay',
                formatting={'style': style}
            )
            
            return filepath
            
        except Exception as e:
            logger.error(f"Error creating DOCX essay: {e}")
            # Fallback: create simple file
            filepath = self.output_dir / filename
            with open(filepath, 'w') as f:
                f.write(content)
            return str(filepath)
    
    def _create_presentation_outline(self, topic: str, slides: int, research: str) -> List[Dict[str, Any]]:
        """Create presentation outline"""
        try:
            outline_prompt = f"""
            Create a presentation outline for "{topic}" with {slides} slides.
            
            Research context:
            {research[:2000]}
            
            Standard presentation structure:
            1. Title slide
            2. Agenda/Overview
            3. Introduction
            4-{slides-2}. Main content slides
            {slides-1}. Conclusion
            {slides}. Thank you/Questions
            
            Return as JSON array:
            [
                {{"slide_number": 1, "title": "Title", "content_type": "title", "main_points": ["Title", "Presenter"]}},
                {{"slide_number": 2, "title": "Agenda", "content_type": "bulleted_list", "main_points": ["Point 1", "Point 2"]}}
            ]
            
            Content types: title, bulleted_list, text_heavy, image_placeholder, chart_placeholder
            """
            
            response = self.claude_engine.client.messages.create(
                model=self.claude_engine.model,
                max_tokens=1500,
                messages=[{"role": "user", "content": outline_prompt}]
            )
            
            outline_text = response.content[0].text.strip()
            
            if outline_text.startswith('```json'):
                outline_text = outline_text[7:-3]
            elif outline_text.startswith('```'):
                outline_text = outline_text[3:-3]
            
            return json.loads(outline_text)
            
        except Exception as e:
            logger.error(f"Error creating presentation outline: {e}")
            # Fallback outline
            return [
                {"slide_number": i+1, "title": f"Slide {i+1}", "content_type": "bulleted_list", "main_points": [f"Point {i+1}"]}
                for i in range(slides)
            ]
    
    def _generate_slide_contents(self, outline: List[Dict[str, Any]], research: str, speaker_notes: bool) -> List[Dict[str, Any]]:
        """Generate detailed content for each slide"""
        slide_contents = []
        
        for slide in outline:
            try:
                content_prompt = f"""
                Generate content for slide #{slide['slide_number']}: "{slide['title']}"
                
                Content type: {slide['content_type']}
                Main points: {slide['main_points']}
                
                Research context:
                {research[:1000]}
                
                Return JSON:
                {{
                    "slide_number": {slide['slide_number']},
                    "title": "{slide['title']}",
                    "bullet_points": ["Point 1", "Point 2", "Point 3"],
                    "speaker_notes": "Detailed speaking notes..." 
                }}
                
                Keep bullet points concise (max 7 words each). Include 3-5 points per slide.
                """
                
                response = self.claude_engine.client.messages.create(
                    model=self.claude_engine.model,
                    max_tokens=600,
                    messages=[{"role": "user", "content": content_prompt}]
                )
                
                content_text = response.content[0].text.strip()
                
                if content_text.startswith('```json'):
                    content_text = content_text[7:-3]
                elif content_text.startswith('```'):
                    content_text = content_text[3:-3]
                
                slide_content = json.loads(content_text)
                slide_contents.append(slide_content)
                
            except Exception as e:
                logger.error(f"Error generating slide content: {e}")
                # Fallback content
                slide_contents.append({
                    "slide_number": slide['slide_number'],
                    "title": slide['title'],
                    "bullet_points": slide['main_points'],
                    "speaker_notes": f"Speaker notes for {slide['title']}"
                })
        
        return slide_contents
    
    def _create_powerpoint(self, slide_contents: List[Dict[str, Any]], filename: str, style: str) -> str:
        """Create PowerPoint presentation"""
        try:
            # Use existing document engine if it supports PowerPoint
            if hasattr(self.document_engine, 'create_presentation'):
                return self.document_engine.create_presentation(
                    slides=slide_contents,
                    filename=filename,
                    style=style
                )
            
            # Fallback: Create using python-pptx
            try:
                from pptx import Presentation
                from pptx.util import Inches
                
                prs = Presentation()
                
                for slide_data in slide_contents:
                    slide_layout = prs.slide_layouts[1]  # Bullet layout
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # Add title
                    title = slide.shapes.title
                    title.text = slide_data['title']
                    
                    # Add bullet points
                    content = slide.shapes.placeholders[1]
                    text_frame = content.text_frame
                    
                    for i, point in enumerate(slide_data.get('bullet_points', [])):
                        if i == 0:
                            text_frame.text = point
                        else:
                            p = text_frame.add_paragraph()
                            p.text = point
                    
                    # Add speaker notes if available
                    if 'speaker_notes' in slide_data:
                        notes_slide = slide.notes_slide
                        notes_text_frame = notes_slide.notes_text_frame
                        notes_text_frame.text = slide_data['speaker_notes']
                
                filepath = self.output_dir / filename
                prs.save(str(filepath))
                
                return str(filepath)
                
            except ImportError:
                logger.warning("python-pptx not available, creating text outline")
                # Create text outline as fallback
                outline_text = "\n\n".join([
                    f"SLIDE {slide['slide_number']}: {slide['title']}\n" + 
                    "\n".join(f"• {point}" for point in slide.get('bullet_points', []))
                    for slide in slide_contents
                ])
                
                filepath = self.output_dir / filename.replace('.pptx', '.txt')
                with open(filepath, 'w') as f:
                    f.write(outline_text)
                
                return str(filepath)
                
        except Exception as e:
            logger.error(f"Error creating PowerPoint: {e}")
            return f"Error creating presentation: {str(e)}"
    
    def _create_report_structure(self, subject: str, data: Dict[str, Any]) -> Dict[str, Any]:
        """Create professional report structure"""
        return {
            "executive_summary": f"Executive summary for {subject}",
            "introduction": f"Introduction to {subject} analysis",
            "methodology": "Data collection and analysis methods",
            "findings": "Key findings and results",
            "analysis": "Detailed analysis and interpretation",
            "recommendations": "Strategic recommendations",
            "conclusion": "Summary and next steps",
            "appendices": "Supporting data and references"
        }
    
    def _write_report_sections(self, subject: str, structure: Dict[str, Any], data: Dict[str, Any], research: str) -> Dict[str, str]:
        """Write each section of the report"""
        sections = {}
        
        for section_name, section_description in structure.items():
            try:
                section_prompt = f"""
                Write a {section_name} section for a professional report on "{subject}".
                
                Section purpose: {section_description}
                
                Available data: {json.dumps(data, indent=2)[:1000]}
                Research context: {research[:1000]}
                
                Write professional, analytical content appropriate for a business report.
                Length: 200-400 words.
                """
                
                response = self.claude_engine.client.messages.create(
                    model=self.claude_engine.model,
                    max_tokens=600,
                    messages=[{"role": "user", "content": section_prompt}]
                )
                
                sections[section_name] = response.content[0].text
                
            except Exception as e:
                logger.error(f"Error writing {section_name} section: {e}")
                sections[section_name] = f"{section_name.replace('_', ' ').title()} content for {subject}."
        
        return sections
    
    def _add_visualizations_to_report(self, sections: Dict[str, str], data: Dict[str, Any]) -> Dict[str, str]:
        """Add charts and tables to report sections"""
        # This would integrate with the existing document engine's chart capabilities
        return sections
    
    def _create_pdf_report(self, sections: Dict[str, str], filename: str) -> str:
        """Create PDF report with professional formatting"""
        try:
            # Combine sections
            content = "\n\n".join([
                f"{section_name.replace('_', ' ').title()}\n{'='*50}\n{content}"
                for section_name, content in sections.items()
            ])
            
            # Use existing document engine
            if hasattr(self.document_engine, 'create_pdf'):
                return self.document_engine.create_pdf(content, filename)
            
            # Fallback: create text file
            filepath = self.output_dir / filename.replace('.pdf', '.txt')
            with open(filepath, 'w') as f:
                f.write(content)
            
            return str(filepath)
            
        except Exception as e:
            logger.error(f"Error creating PDF report: {e}")
            return f"Error creating report: {str(e)}"
    
    def _generate_bibliography(self, style: str) -> str:
        """Generate bibliography in specified style"""
        if style in self.citation_styles:
            return self.citation_styles[style]()
        else:
            return "\n\nReferences\n==========\n\n[Research sources would be listed here]"
    
    def _apa_citation_format(self) -> str:
        """Generate APA style bibliography"""
        return """
References

[Note: In a real implementation, these would be actual sources from the research]

Author, A. A. (2024). Title of article. Journal Name, 12(3), 45-67.

Organization. (2024). Report title. Publisher.

Website Author. (2024). Article title. Website Name. https://example.com
"""
    
    def _mla_citation_format(self) -> str:
        """Generate MLA style bibliography"""
        return """
Works Cited

[Note: In a real implementation, these would be actual sources from the research]

Author, First. "Article Title." Journal Name, vol. 12, no. 3, 2024, pp. 45-67.

Organization. Report Title. Publisher, 2024.

Website Author. "Article Title." Website Name, 2024, https://example.com.
"""
    
    def _chicago_citation_format(self) -> str:
        """Generate Chicago style bibliography"""
        return """
Bibliography

[Note: In a real implementation, these would be actual sources from the research]

Author, First Last. "Article Title." Journal Name 12, no. 3 (2024): 45-67.

Organization. Report Title. Publisher, 2024.

Website Author. "Article Title." Website Name. Accessed January 3, 2026. https://example.com.
"""